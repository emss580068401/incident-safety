import os
import sys
import subprocess

try:
    import pypdf
except ImportError:
    subprocess.check_call([sys.executable, "-m", "pip", "install", "pypdf"])
    import pypdf

try:
    import pptx
except ImportError:
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    import pptx

files = [
    r'd:\Antigravity\AI專用\法規教材\消防職安及法規專章介紹-消防署綜企組 陳坤佐組長.pdf',
    r'd:\Antigravity\AI專用\法規教材\消防三權(2小時)-吳俊螢.pdf',
    r'd:\Antigravity\AI專用\事故安全官\115年3月份事故安全官分享(較新).pdf',
    r'd:\Antigravity\AI專用\事故安全官\114事故安全官實體課程簡報.pptx',
    r'd:\Antigravity\AI專用\救災安全管制及後勤照護\救災安全管制及後勤照護-4小時.pdf'
]

output = []

for filepath in files:
    output.append(f"\n===========================\nFile: {os.path.basename(filepath)}\n===========================\n")
    if not os.path.exists(filepath):
        output.append("File not found.")
        continue
    
    if filepath.endswith('.pdf'):
        try:
            reader = pypdf.PdfReader(filepath)
            # Try to get headings or just extract first few lines of each page
            for i, page in enumerate(reader.pages):
                text = page.extract_text()
                if text:
                    lines = text.split('\n')
                    # Just keep the first 3 lines of each page to understand topics
                    title = " ".join(lines[:3][:50])
                    if title.strip():
                        output.append(f"Page {i+1}: {title[:100]}")
        except Exception as e:
            output.append(f"Error reading PDF: {e}")
            
    elif filepath.endswith('.pptx'):
        try:
            prs = pptx.Presentation(filepath)
            for i, slide in enumerate(prs.slides):
                titles = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        titles.append(shape.text.replace('\n', ' ').strip()[:50])
                if titles:
                    output.append(f"Slide {i+1}: {titles[0]}")
        except Exception as e:
            output.append(f"Error reading PPTX: {e}")

with open("content_outline.txt", "w", encoding="utf-8") as f:
    f.write("\n".join(output))

print("Outline extracted to content_outline.txt")
